from flask import Flask, render_template, request
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.asymmetric import dh
from cryptography.hazmat.backends import default_backend
import os
import docx
from pptx import Presentation

app = Flask(__name__)

# Generate Diffie-Hellman key exchange
def generate_key():
    parameters = dh.generate_parameters(generator=2, key_size=512, backend=default_backend())
    private_key = parameters.generate_private_key()
    public_key = private_key.public_key()
    return private_key, public_key

def encrypt_text(private_key, text):
    salt = os.urandom(16)
    kdf = PBKDF2HMAC(
        algorithm=hashes.SHA256(),
        length=32,
        salt=salt,
        iterations=100000,
        backend=default_backend()
    )
    key = kdf.derive(private_key.private_numbers().x.to_bytes(64, 'big'))
    # Simple XOR encryption
    encrypted_text = ''.join(chr(ord(c) ^ key[0]) for c in text)
    return encrypted_text, key

@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        document_file = request.files['document']
        document_path = os.path.join("uploads", document_file.filename)
        document_file.save(document_path)

        # Read and encrypt the document text
        if document_file.filename.endswith('.docx'):
            doc = docx.Document(document_path)
            document_text = '\n'.join([para.text for para in doc.paragraphs])
        elif document_file.filename.endswith('.pptx'):
            prs = Presentation(document_path)
            document_text = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        document_text += shape.text + '\n'
        else:
            return "Unsupported file type!"

        private_key, _ = generate_key()
        encrypted_text, key = encrypt_text(private_key, document_text)

        # For demonstration, just display the encrypted text on the page
        return f'''
        <html>
        <body>
        <h1>Files Processed Successfully!</h1>
        <h2>Encrypted Content:</h2>
        <p>{encrypted_text}</p>
        <a href="/">Upload another document</a>
        </body>
        </html>
        '''

    return '''
    <html>
    <body>
    <h1>Upload Document to Encrypt</h1>
    <form action="/" method="POST" enctype="multipart/form-data">
        <input type="file" name="document" accept=".docx,.pptx"><br><br>
        <input type="submit" value="Encrypt and Download">
    </form>
    </body>
    </html>
    '''

if __name__ == '__main__':
    os.makedirs("uploads", exist_ok=True)
    app.run(debug=True)
