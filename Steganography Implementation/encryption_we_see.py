from flask import Flask, render_template, request, send_file
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.asymmetric import dh
import os
from docx import Document
from pptx import Presentation
import io
import pyexiv2
from PIL import Image
from PIL import PngImagePlugin




app = Flask(__name__)

# Create the uploads directory if it doesn't exist
if not os.path.exists('uploads'):
    os.makedirs('uploads')

# Generate Diffie-Hellman key exchange
def generate_key():
    parameters = dh.generate_parameters(generator=2, key_size=512)
    private_key = parameters.generate_private_key()
    public_key = private_key.public_key()
    return private_key, public_key

def encrypt_text(private_key, text):
    salt = os.urandom(16)
    kdf = PBKDF2HMAC(
        algorithm=hashes.SHA256(),
        length=32,
        salt=salt,
        iterations=100000
    )
    key = kdf.derive(private_key.private_numbers().x.to_bytes(64, 'big'))
    # Simple XOR encryption
    encrypted_text = ''.join(chr(ord(c) ^ key[0]) for c in text)
    return encrypted_text, key

# Store key in image metadata
def store_key_in_image(image_path, key):
    image = Image.open(image_path)
    metadata = PngImagePlugin.PngInfo()
    metadata.add_text("Software", key.hex())
    
    # Save the image with the new metadata
    image.save(image_path, "PNG", pnginfo=metadata)

@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        document_file = request.files['document']
        document_path = os.path.join("uploads", document_file.filename)
        document_file.save(document_path)

        # Check the file extension to process accordingly
        if document_file.filename.endswith('.docx'):
            # Process Word document
            document = Document(document_path)
            text_content = []
            images = []

            for para in document.paragraphs:
                text_content.append(para.text)

            for rel in document.part.rels.values():
                if "image" in rel.target_ref:
                    image_data = rel.target_part.blob
                    image_name = rel.target_ref.split("/")[-1]
                    images.append((image_name, image_data))

            text = "\n".join(text_content)

        elif document_file.filename.endswith('.pptx'):
            # Process PowerPoint presentation
            presentation = Presentation(document_path)
            text_content = []
            images = []

            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
                    if shape.shape_type == 13:  # Shape type 13 is a picture
                        image_stream = io.BytesIO()
                        shape.image.save(image_stream)
                        images.append((shape.image.filename, image_stream.getvalue()))

            text = "\n".join(text_content)

        # Encrypt text using Diffie-Hellman
        private_key, _ = generate_key()
        encrypted_text, key = encrypt_text(private_key, text)

        # Process images
        for image_name, image_data in images:
            image_path = os.path.join("uploads", image_name)
            with open(image_path, 'wb') as img_file:
                img_file.write(image_data)
            store_key_in_image(image_path, key)

        return send_file(document_path, as_attachment=True)

    return '''
    <html>
    <body>
    <h1>Upload Document (Word or PowerPoint) to Encrypt</h1>
    <form action="/" method="POST" enctype="multipart/form-data">
        <input type="file" name="document" accept=".docx,.pptx"><br><br>
        <input type="submit" value="Encrypt and Download">
    </form>
    </body>
    </html>
    '''

if __name__ == '__main__':
    app.run(debug=True)
